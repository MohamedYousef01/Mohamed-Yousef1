import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import cast
from pptx.shapes.autoshape import Shape

# ── Colours ─────────────────────────────────────────────────────────────────
_HEADER_BG  = RGBColor(0x1E, 0x40, 0x6E)   # dark navy
_HEADER_FG  = RGBColor(0xFF, 0xFF, 0xFF)   # white
_ROW_EVEN   = RGBColor(0xEF, 0xF4, 0xFB)   # light blue-grey
_ROW_ODD    = RGBColor(0xFF, 0xFF, 0xFF)   # white
_TEXT       = RGBColor(0x1F, 0x29, 0x37)   # near-black
_TITLE_CLR  = RGBColor(0x1F, 0x29, 0x37)


def _set_cell(cell, text, bold=False, font_size=12,
              fg=_TEXT, bg=None, align=PP_ALIGN.LEFT):
    """Write text into a cell safely — works even when text is empty."""
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()

    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = fg


def _add_bullet_slide(prs, slide_data):
    layout = prs.slide_layouts[1]
    slide_obj = prs.slides.add_slide(layout)

    if slide_obj.shapes.title:
        slide_obj.shapes.title.text = slide_data.get("title", "No Title")

    content = None
    for shape in slide_obj.shapes:
        if shape.has_text_frame and shape != slide_obj.shapes.title:
            content = cast(Shape, shape)
            break

    if content:
        tf = content.text_frame
        tf.clear()
        for i, bullet in enumerate(slide_data.get("bullets", [])):
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0


def _add_table_slide(prs, slide_data):
    # Use blank layout — find it by name, fall back to last layout
    blank = next(
        (l for l in prs.slide_layouts if l.name.lower() == "blank"),
        prs.slide_layouts[-1]
    )
    slide_obj = prs.slides.add_slide(blank)

    # Remove every inherited placeholder so nothing shows "Click to add..."
    for ph in list(slide_obj.placeholders):
        ph._element.getparent().remove(ph._element)

    table_data = slide_data.get("table", {})
    headers    = table_data.get("headers", [])
    rows       = table_data.get("rows", [])

    if not headers:
        return

    # ── Title ────────────────────────────────────────────────────────────────
    title_box = slide_obj.shapes.add_textbox(
        Inches(0.45), Inches(0.25), Inches(9.1), Inches(0.65)
    )
    tf = title_box.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    run  = para.add_run()
    run.text            = slide_data.get("title", "Comparison")
    run.font.size       = Pt(26)
    run.font.bold       = True
    run.font.color.rgb  = _TITLE_CLR

    # ── Table geometry ────────────────────────────────────────────────────────
    n_cols      = len(headers)
    n_rows      = len(rows) + 1        # header + data rows
    ROW_H       = Inches(0.52)         # fixed per-row height
    tbl_left    = Inches(0.45)
    tbl_top     = Inches(1.05)
    tbl_width   = Inches(9.1)
    tbl_height  = ROW_H * n_rows

    tbl_shape = slide_obj.shapes.add_table(
        n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height
    )
    tbl = tbl_shape.table

    # Distribute columns evenly
    col_width = tbl_width // n_cols
    for col in tbl.columns:
        col.width = col_width

    # Fix every row to the same height
    for row in tbl.rows:
        row.height = ROW_H

    # ── Header row ────────────────────────────────────────────────────────────
    for col_idx, header in enumerate(headers):
        _set_cell(
            tbl.cell(0, col_idx), header,
            bold=True, font_size=13,
            fg=_HEADER_FG, bg=_HEADER_BG,
            align=PP_ALIGN.CENTER
        )

    # ── Data rows ─────────────────────────────────────────────────────────────
    for row_idx, row in enumerate(rows):
        bg = _ROW_EVEN if row_idx % 2 == 0 else _ROW_ODD
        for col_idx in range(n_cols):
            value = row[col_idx] if col_idx < len(row) else ""
            _set_cell(
                tbl.cell(row_idx + 1, col_idx), value,
                font_size=12, fg=_TEXT, bg=bg
            )


def create_ppt(slides_json):
    if isinstance(slides_json, str):
        slides_json = json.loads(slides_json)

    prs = Presentation()

    for slide in slides_json.get("slides", []):
        if "table" in slide:
            _add_table_slide(prs, slide)
        else:
            _add_bullet_slide(prs, slide)

    file_path = "presentation_output.pptx"
    prs.save(file_path)
    return file_path
